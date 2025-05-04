# pptx_server.py
import io
import os
import subprocess
import sys # Import sys for stderr printing
import tempfile
import shutil # For finding soffice
import base64 # <-- Add this import
from pathlib import Path
from typing import Optional, List, Union, Tuple # <-- Add Tuple

import pptx # Import the module directly
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor # <-- Add RGBColor

from fastmcp import FastMCP, Image, Context
from fastmcp.resources import FileResource

# --- Configuration ---
# Directory to store generated presentations
SAVE_DIR = Path("./presentations")
SAVE_DIR.mkdir(exist_ok=True)

# --- !! IMPORTANT !! LibreOffice Dependency ---
# Set this if 'soffice' is not in your system PATH
SOFFICE_PATH = None # e.g., "/usr/bin/soffice" or "C:\\Program Files\\LibreOffice\\program\\soffice.bin"
# --- End Configuration ---


# --- FastMCP Server Definition ---
# NOTE: python-pptx is included, but LibreOffice CANNOT be installed via pip.
# Deployment requires manual installation of LibreOffice on the host system.
#
# Server host/port can be set via HOST and PORT environment variables, defaulting to 127.0.0.1:8000
HOST = os.environ.get("HOST", "127.0.0.1")
PORT = int(os.environ.get("PORT", "8000"))

mcp = FastMCP(
    "PowerPoint Creator 📊 (with Image Rendering)",
    dependencies=["python-pptx", "pillow"] # Added pillow dependency for Image helper
)

# --- Helper Functions ---

def _find_soffice() -> str:
    """Finds the LibreOffice executable."""
    if SOFFICE_PATH and Path(SOFFICE_PATH).exists():
        return SOFFICE_PATH
    
    soffice_cmd = "soffice"
    if os.name == 'nt': # Windows
        soffice_cmd = "soffice.exe" # Check PATH
        # Common Windows install locations (add more if needed)
        program_files = os.environ.get("ProgramFiles", "C:\\Program Files")
        program_files_x86 = os.environ.get("ProgramFiles(x86)", "C:\\Program Files (x86)")
        possible_paths = [
            Path(program_files) / "LibreOffice" / "program" / "soffice.exe",
            Path(program_files_x86) / "LibreOffice" / "program" / "soffice.exe",
        ]
        for p in possible_paths:
            if p.exists():
                return str(p)

    # For Linux/macOS, shutil.which checks the PATH
    found_path = shutil.which(soffice_cmd)
    if found_path:
        return found_path

    raise RuntimeError(
        "LibreOffice 'soffice' executable not found in PATH or configured SOFFICE_PATH. "
        "Image rendering requires LibreOffice installation."
    )


def _get_presentation_path(filename: str) -> Path:
    """Constructs the full path for the presentation file."""
    if not filename.endswith(".pptx"):
        filename += ".pptx"
    # Basic security: ensure filename doesn't escape SAVE_DIR
    path = SAVE_DIR / Path(filename).name
    if not path.resolve().is_relative_to(SAVE_DIR.resolve()):
         raise ValueError("Invalid filename causing path traversal.")
    return path

def _load_presentation(filename: str) -> pptx.Presentation: # Use qualified name
    """Loads a presentation or creates a new one if it doesn't exist."""
    path = _get_presentation_path(filename)
    if path.exists():
        try:
            return pptx.Presentation(str(path)) # Use qualified name
        except Exception as e:
            # Catch other potential errors during loading (e.g., file corruption)
            raise ValueError(f"Error loading presentation '{filename}': {e}")
    else:
        return pptx.Presentation() # Use qualified name

def _save_presentation(prs: pptx.Presentation, filename: str): # Use qualified name
    """Saves the presentation object to the specified file."""
    path = _get_presentation_path(filename)
    try:
        prs.save(path)
    except Exception as e:
        raise IOError(f"Error saving presentation to '{path}': {e}")


def _get_slide(prs: pptx.Presentation, slide_index: int): # Use qualified name
    """Gets a specific slide by index, raising user-friendly errors."""
    if not isinstance(slide_index, int) or slide_index < 0:
        raise ValueError(f"Slide index must be a non-negative integer, got {slide_index}.")
    try:
        # Presentation.slides is list-like, check bounds explicitly
        if slide_index >= len(prs.slides):
             raise IndexError # Trigger the same error path
        return prs.slides[slide_index]
    except IndexError:
        raise ValueError(f"Invalid slide index {slide_index}. Presentation has {len(prs.slides)} slides (0-indexed).")

def _parse_shape_type(shape_name: str) -> MSO_SHAPE:
    """Converts a string shape name to an MSO_SHAPE enum."""
    try:
        # Convert to upper case for case-insensitive matching
        return getattr(MSO_SHAPE, shape_name.upper())
    except AttributeError:
        # Provide a list of common/useful shapes in the error
        common_shapes = ["RECTANGLE", "OVAL", "ROUNDED_RECTANGLE", "DIAMOND",
                         "ISOSCELES_TRIANGLE", "RIGHT_ARROW", "LEFT_ARROW",
                         "UP_ARROW", "DOWN_ARROW", "PENTAGON", "HEXAGON",
                         "CHEVRON", "STAR_5_POINT", "FLOWCHART_PROCESS",
                         "FLOWCHART_DECISION", "FLOWCHART_TERMINATOR",
                         "FLOWCHART_DATA", "LINE_CALLOUT_1"] # etc.
        raise ValueError(f"Unknown shape type '{shape_name}'. Try one of: {', '.join(common_shapes)}...")

def _get_shape_by_id(slide: pptx.slide.Slide, shape_id: int):
    """Finds a shape on the slide by its unique ID."""
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    raise ValueError(f"Shape with ID {shape_id} not found on slide {slide.slide_id}. Available IDs: {[s.shape_id for s in slide.shapes]}")

# --- MCP Tools (Create, Add Slide, Add Elements - same as before) ---

@mcp.tool()
def create_or_clear_presentation(filename: str) -> str:
    """
    Creates a new, empty presentation with the given filename,
    or clears an existing one. Overwrites if the file exists.
    """
    # Ensure the presentation path is valid before creating
    pptx_path = _get_presentation_path(filename)
    prs = pptx.Presentation() # Use qualified name
    _save_presentation(prs, filename)
    return f"Presentation '{filename}' created/cleared successfully in '{SAVE_DIR}'."

@mcp.tool()
def add_slide(filename: str, layout_index: int = 5) -> str:
    """
    Adds a new slide to the presentation using a specified layout index.
    Common layouts: 0=Title, 1=Title+Content, 5=Title Only, 6=Blank.
    Returns the index of the newly added slide.
    """
    prs = _load_presentation(filename)
    if not (0 <= layout_index < len(prs.slide_layouts)):
        raise ValueError(f"Invalid layout_index {layout_index}. Must be between 0 and {len(prs.slide_layouts) - 1}.")
    slide_layout = prs.slide_layouts[layout_index]
    prs.slides.add_slide(slide_layout)
    new_slide_index = len(prs.slides) - 1 # Index is 0-based
    _save_presentation(prs, filename)
    return f"Added slide {new_slide_index} with layout {layout_index} to '{filename}'. New slide count: {len(prs.slides)}."


@mcp.tool()
def add_title_and_content(filename: str, slide_index: int, title: str, content: str) -> str:
    """
    Adds text to the title and main content placeholder of a specific slide.
    Assumes the slide layout has these placeholders (e.g., layout index 1).
    """
    prs = _load_presentation(filename)
    slide = _get_slide(prs, slide_index)

    title_shape = None
    content_placeholder = None

    # Find title placeholder (usually idx 0 or specific name)
    if slide.shapes.title:
        title_shape = slide.shapes.title
    else:
        for shape in slide.placeholders:
            if shape.name.lower().startswith("title"):
                title_shape = shape
                break

    if title_shape:
        title_shape.text = title
    else:
         print(f"Warning: Slide {slide_index} does not have a standard title placeholder.")


    # Find the main content placeholder (often index 1, but search robustly)
    for shape in slide.placeholders:
        # Check common indices or names
        if shape.placeholder_format.idx == 1 or \
           shape.name.lower().startswith("content placeholder") or \
           shape.name.lower().startswith("text placeholder") or \
           shape.name.lower().startswith("body"):

           # Avoid assigning the title shape if it was also found this way
           if shape != title_shape:
                content_placeholder = shape
                break

    if not content_placeholder and len(slide.placeholders) > 1 and slide.placeholders[1] != title_shape:
         # Fallback to index 1 if different from title
         content_placeholder = slide.placeholders[1]

    if content_placeholder:
        tf = content_placeholder.text_frame
        tf.text = content # Set first paragraph
        # Optionally clear other paragraphs if needed: while len(tf.paragraphs) > 1: tf._remove_paragraph(tf.paragraphs[-1])
        # Optionally add more paragraphs for bullet points if content has newlines etc.
    else:
        print(f"Warning: Slide {slide_index} does not seem to have a standard content placeholder.")
        # As a last resort, could add a new textbox, but maybe better to inform the user.

    _save_presentation(prs, filename)
    return f"Attempted to add title and content to slide {slide_index} in '{filename}'."


@mcp.tool()
def add_textbox(
    filename: str,
    slide_index: int,
    text: str,
    left_inches: float,
    top_inches: float,
    width_inches: float,
    height_inches: float,
    font_size_pt: int = 0, # Changed: Use 0 to indicate 'not set' instead of Optional[int]
    bold: bool = False
) -> str:
    """
    Adds a textbox with specified text, position, and dimensions (in inches) to a slide.
    Set font_size_pt to 0 or less to use the default font size.
    """ # Updated docstring
    prs = _load_presentation(filename)
    slide = _get_slide(prs, slide_index)
    left, top = Inches(left_inches), Inches(top_inches)
    width, height = Inches(width_inches), Inches(height_inches)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    # Handle multi-line text properly
    tf.text = text.split('\n')[0] # First line
    for line in text.split('\n')[1:]:
        p = tf.add_paragraph()
        p.text = line

    # Apply formatting to all paragraphs in the textbox
    for p in tf.paragraphs:
        p.font.bold = bold
        # Changed condition: Check for > 0 instead of checking for None
        if font_size_pt > 0:
            p.font.size = Pt(font_size_pt)

    _save_presentation(prs, filename)
    return f"Added textbox to slide {slide_index} in '{filename}'."


@mcp.tool()
def add_shape(
    filename: str,
    slide_index: int,
    shape_type_name: str,
    left_inches: float,
    top_inches: float,
    width_inches: float,
    height_inches: float,
    text: Optional[str] = None,
) -> str:
    """
    Adds an AutoShape (like RECTANGLE, OVAL, FLOWCHART_PROCESS) to a slide.
    Specify position and dimensions in inches. Optionally add text to the shape.
    Returns a confirmation message including the unique ID of the created shape.
    """ # Updated docstring
    prs = _load_presentation(filename)
    slide = _get_slide(prs, slide_index)
    shape_enum = _parse_shape_type(shape_type_name)
    left, top = Inches(left_inches), Inches(top_inches)
    width, height = Inches(width_inches), Inches(height_inches)

    shape = slide.shapes.add_shape(shape_enum, left, top, width, height)
    shape_id = shape.shape_id # Get the unique ID

    if text:
        tf = shape.text_frame
        # Handle multi-line text in shapes too
        tf.text = text.split('\n')[0]
        for line in text.split('\n')[1:]:
             p = tf.add_paragraph()
             p.text = line
        tf.word_wrap = True # Enable word wrap within the shape

    _save_presentation(prs, filename)
    # Include the shape_id in the return message
    return f"Added shape '{shape_type_name}' (ID: {shape_id}) to slide {slide_index} in '{filename}'."

@mcp.tool()
def add_connector(
    filename: str,
    slide_index: int,
    start_shape_id: int,
    end_shape_id: int,
    connector_type_name: str = "ELBOW", # Common default: STRAIGHT, ELBOW, CURVE
    start_connection_point_idx: int = 3, # Default: Mid-right side for many shapes
    end_connection_point_idx: int = 1,   # Default: Mid-left side for many shapes
) -> str:
    """
    Adds a connector shape between two existing shapes identified by their IDs.
    Defaults to an ELBOW connector from the right side of the start shape
    to the left side of the end shape.

    Args:
        filename: The presentation filename.
        slide_index: The 0-based index of the slide.
        start_shape_id: The unique ID of the shape where the connector starts.
        end_shape_id: The unique ID of the shape where the connector ends.
        connector_type_name: Type of connector (e.g., "STRAIGHT", "ELBOW", "CURVE").
        start_connection_point_idx: Index of the connection point on the start shape (0=center, 1-N=perimeter points).
        end_connection_point_idx: Index of the connection point on the end shape.

    Returns:
        Confirmation message including the connector's shape ID.
    """
    prs = _load_presentation(filename)
    slide = _get_slide(prs, slide_index)

    # Find the start and end shapes
    start_shape = _get_shape_by_id(slide, start_shape_id)
    end_shape = _get_shape_by_id(slide, end_shape_id)

    # Parse connector type
    try:
        connector_enum = getattr(MSO_CONNECTOR, connector_type_name.upper())
    except AttributeError:
        raise ValueError(f"Unknown connector type '{connector_type_name}'. Try: STRAIGHT, ELBOW, CURVE.")

    # Add the connector shape (initial position doesn't matter much)
    connector = slide.shapes.add_connector(
        connector_enum, Inches(1), Inches(1), Inches(2), Inches(2) # Arbitrary start/end
    )

    # Connect the shapes
    try:
        connector.begin_connect(start_shape, start_connection_point_idx)
    except Exception as e:
        # Provide more context on error
        print(f"Warning: Could not connect start of connector to shape {start_shape_id} at point {start_connection_point_idx}. Error: {e}. Ensure the connection point index is valid for the shape type.")
        # Attempt to connect to center as fallback?
        try:
            print("Attempting fallback connection to center (point 0) for start shape.")
            connector.begin_connect(start_shape, 0)
        except Exception as e2:
             print(f"Fallback connection to center also failed: {e2}")
             # Proceed without connection if fallback fails, user might fix later

    try:
        connector.end_connect(end_shape, end_connection_point_idx)
    except Exception as e:
        print(f"Warning: Could not connect end of connector to shape {end_shape_id} at point {end_connection_point_idx}. Error: {e}. Ensure the connection point index is valid for the shape type.")
        # Attempt to connect to center as fallback?
        try:
            print("Attempting fallback connection to center (point 0) for end shape.")
            connector.end_connect(end_shape, 0)
        except Exception as e2:
            print(f"Fallback connection to center also failed: {e2}")
            # Proceed without connection if fallback fails

    _save_presentation(prs, filename)
    connector_id = connector.shape_id
    return f"Added {connector_type_name} connector (ID: {connector_id}) from shape {start_shape_id} (point {start_connection_point_idx}) to shape {end_shape_id} (point {end_connection_point_idx}) on slide {slide_index}."

@mcp.tool()
def delete_shape(filename: str, slide_index: int, shape_id: int) -> str:
    """
    Deletes a specific shape from a slide using its unique ID.

    Args:
        filename: The presentation filename.
        slide_index: The 0-based index of the slide.
        shape_id: The unique ID of the shape to delete.

    Returns:
        Confirmation message.
    """
    prs = _load_presentation(filename)
    slide = _get_slide(prs, slide_index)

    # Find the shape to delete
    shape_to_delete = _get_shape_by_id(slide, shape_id)

    # Remove the shape element from its parent
    sp = shape_to_delete._element # Access the underlying XML element
    sp.getparent().remove(sp)   # Remove the element from the shapes collection

    _save_presentation(prs, filename)
    return f"Deleted shape with ID {shape_id} from slide {slide_index} in '{filename}'."

@mcp.tool()
def modify_shape(
    filename: str,
    slide_index: int,
    shape_id: int,
    text: Optional[str] = None,
    left_inches: Optional[float] = None,
    top_inches: Optional[float] = None,
    width_inches: Optional[float] = None,
    height_inches: Optional[float] = None,
    font_size_pt: Optional[int] = None,
    bold: Optional[bool] = None,
    fill_color_rgb: Optional[Tuple[int, int, int]] = None,
    line_color_rgb: Optional[Tuple[int, int, int]] = None,
    line_width_pt: Optional[float] = None
) -> str:
    """
    Modifies properties of an existing shape identified by its ID.
    Allows changing text, position, size, font attributes, fill color, and line style.
    Only provided parameters are changed.

    Args:
        filename: The presentation filename.
        slide_index: The 0-based index of the slide.
        shape_id: The unique ID of the shape to modify.
        text: New text content for the shape (replaces existing text).
        left_inches: New horizontal position from the left edge (in inches).
        top_inches: New vertical position from the top edge (in inches).
        width_inches: New width (in inches).
        height_inches: New height (in inches).
        font_size_pt: New font size (in points) for all text in the shape.
        bold: Set font bold state (True or False) for all text in the shape.
        fill_color_rgb: Tuple of (R, G, B) values (0-255) for solid fill color.
        line_color_rgb: Tuple of (R, G, B) values (0-255) for line color.
        line_width_pt: New line width (in points).

    Returns:
        Confirmation message summarizing the changes made.
    """
    prs = _load_presentation(filename)
    slide = _get_slide(prs, slide_index)
    shape = _get_shape_by_id(slide, shape_id)
    changes_made = []

    # Modify Position/Size
    if left_inches is not None:
        shape.left = Inches(left_inches)
        changes_made.append("position (left)")
    if top_inches is not None:
        shape.top = Inches(top_inches)
        changes_made.append("position (top)")
    if width_inches is not None:
        shape.width = Inches(width_inches)
        changes_made.append("size (width)")
    if height_inches is not None:
        shape.height = Inches(height_inches)
        changes_made.append("size (height)")

    # Modify Text and Font (if shape has text frame)
    if shape.has_text_frame:
        if text is not None:
            tf = shape.text_frame
            tf.clear() # Clear existing paragraphs before adding new text
            # Handle multi-line text properly
            lines = text.split('\n')
            tf.text = lines[0]
            for line in lines[1:]:
                p = tf.add_paragraph()
                p.text = line
            tf.word_wrap = True # Ensure word wrap is enabled
            changes_made.append("text content")

        font_changed = False
        if font_size_pt is not None:
            for p in shape.text_frame.paragraphs:
                p.font.size = Pt(font_size_pt)
            font_changed = True
        if bold is not None:
            for p in shape.text_frame.paragraphs:
                p.font.bold = bold
            font_changed = True
        if font_changed:
             changes_made.append("font attributes (size/bold)")

    # Modify Fill Color
    if fill_color_rgb is not None:
        if len(fill_color_rgb) == 3:
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*fill_color_rgb)
            changes_made.append("fill color")
        else:
            print(f"Warning: Invalid RGB tuple {fill_color_rgb} for fill color. Expected (R, G, B).")

    # Modify Line Style
    line_changed = False
    line = shape.line
    if line_color_rgb is not None:
        if len(line_color_rgb) == 3:
            line.color.rgb = RGBColor(*line_color_rgb)
            line_changed = True
        else:
             print(f"Warning: Invalid RGB tuple {line_color_rgb} for line color. Expected (R, G, B).")
    if line_width_pt is not None:
        line.width = Pt(line_width_pt)
        line_changed = True
    if line_changed:
        changes_made.append("line style (color/width)")

    if not changes_made:
        return f"No valid modifications specified for shape ID {shape_id} on slide {slide_index}."

    _save_presentation(prs, filename)
    return f"Modified shape ID {shape_id} on slide {slide_index}: updated {', '.join(changes_made)}."

@mcp.tool()
def add_picture(
    filename: str,
    slide_index: int,
    image: bytes, # Changed type hint from Image to bytes
    left_inches: float,
    top_inches: float,
    width_inches: Optional[float] = None,
    height_inches: Optional[float] = None,
) -> str:
    """
    Adds a picture to a slide from provided image data.
    Specify position in inches. Optionally specify width OR height in inches to scale.
    If neither width nor height is given, the image's native size is used.
    """
    prs = _load_presentation(filename)
    slide = _get_slide(prs, slide_index)
    left, top = Inches(left_inches), Inches(top_inches)
    width = Inches(width_inches) if width_inches is not None else None
    height = Inches(height_inches) if height_inches is not None else None

    # Use BytesIO to pass image data (which is now bytes) to python-pptx
    image_stream = io.BytesIO(image) # Use image directly as it's now bytes

    slide.shapes.add_picture(image_stream, left, top, width=width, height=height)

    _save_presentation(prs, filename)
    return f"Added picture to slide {slide_index} in '{filename}'."


# --- MCP Resources (Description + NEW Image Rendering) ---

@mcp.resource("pptx://{filename}/slide/{slide_index}/description")
async def get_slide_content_description(filename: str, slide_index: str) -> str:
    """
    Provides a textual description of the shapes and text on a specific slide.
    Useful for the AI to 'confirm' the slide content without actual image rendering.
    [Reliable & Lightweight]
    """
    try:
        idx = int(slide_index)
    except ValueError:
        raise ValueError("Slide index must be an integer.")

    prs = _load_presentation(filename)
    slide = _get_slide(prs, idx) # _get_slide handles index errors

    description = f"--- Slide {idx} Content Description for '{filename}' ---\n"
    description += f"Layout: {slide.slide_layout.name}\n"
    description += f"Number of Shapes: {len(slide.shapes)}\n\n"

    for i, shape in enumerate(slide.shapes):
        shape_type = shape.shape_type
        # Use name attribute if available (e.g., for MSO_SHAPE enums), otherwise use string representation
        type_name = getattr(shape_type, 'name', str(shape_type))
        desc = f"Shape {i}: Type={type_name}"
        # Include the unique shape ID
        desc += f", ID={shape.shape_id}"
        try:
             desc += f", Left={shape.left.inches:.2f}\", Top={shape.top.inches:.2f}\", Width={shape.width.inches:.2f}\", Height={shape.height.inches:.2f}\""
        except AttributeError:
             desc += " (Position/Size not available)" # Handle shapes without these properties if they exist

        if shape.has_text_frame and shape.text.strip():
             # Truncate long text for brevity
            text_preview = (shape.text[:75] + '...') if len(shape.text) > 75 else shape.text
            desc += f", Text='{text_preview.replace(chr(11), ' ').replace('\n', ' ')}'" # Replace VT and newlines for single line desc

        description += desc + "\n"

    if not slide.shapes:
        description += "(Slide is empty)\n"

    description += "--- End Description ---"
    return description


@mcp.tool() # Changed from resource to tool
def get_slide_image(filename: str, slide_index: int) -> Image: # Changed slide_index type hint
    """
    Renders a specific slide as a PNG image using LibreOffice and returns it as an Image object.
    Requires LibreOffice installed and accessible on the server. May be slow.
    Use get_slide_content_description for a faster, text-based check.

    Args:
        filename: The name of the presentation file (e.g., "my_presentation.pptx").
        slide_index: The 0-based index of the slide to render.

    Returns:
        A fastmcp.Image object containing the PNG image data.

    Raises:
        ValueError: If filename is invalid, slide_index is not an integer or out of bounds.
        FileNotFoundError: If the presentation file or the generated PNG is not found.
        RuntimeError: If LibreOffice is not found or the conversion process fails/times out.
    """
    # Function body remains largely the same, just ensure idx is used correctly
    print(f"Attempting slide image rendering for slide {slide_index} of '{filename}' using LibreOffice...", file=sys.stderr)
    print("INFO: This requires LibreOffice to be installed and configured on the server.", file=sys.stderr)

    # Parameter is already slide_index: int, no need for conversion here
    idx = slide_index # Use the integer index directly

    pptx_path = _get_presentation_path(filename)
    if not pptx_path.exists():
        raise ValueError(f"Presentation file '{filename}' not found.")

    # Check number of slides BEFORE conversion to validate index early
    try:
        prs_check = pptx.Presentation(str(pptx_path)) # Use qualified name
        if idx >= len(prs_check.slides):
             raise ValueError(f"Invalid slide index {idx}. Presentation '{filename}' has {len(prs_check.slides)} slides (0-indexed).")
        del prs_check # Close file handle
    except Exception as e:
        print(f"ERROR: Failed to quickly check slide count for {filename}: {e}", file=sys.stderr)
        # Proceed cautiously, LibreOffice might handle corrupted files differently

    soffice = _find_soffice() # Raises RuntimeError if not found

    with tempfile.TemporaryDirectory() as temp_dir:
        temp_dir_path = Path(temp_dir)
        print(f"INFO: Using temporary directory: {temp_dir}", file=sys.stderr)

        # Command to convert the *entire* presentation to PNGs in the temp dir
        # soffice usually names output files based on the input filename + slide number
        cmd = [
            soffice,
            "--headless",          # Run without UI
            "--convert-to", "png", # Convert to PNG format
            "--outdir", str(temp_dir_path), # Output directory
            str(pptx_path)         # Input PPTX file
        ]

        try:
            print(f"INFO: Running LibreOffice command: {' '.join(cmd)}", file=sys.stderr)
            # Use a timeout to prevent hanging indefinitely
            timeout_seconds = 60
            process = subprocess.run(cmd, capture_output=True, text=True, check=False, timeout=timeout_seconds)

            if process.returncode != 0:
                print(f"ERROR: LibreOffice conversion failed! Return code: {process.returncode}", file=sys.stderr)
                print(f"ERROR: Stderr: {process.stderr}", file=sys.stderr)
                print(f"ERROR: Stdout: {process.stdout}", file=sys.stderr)
                raise RuntimeError(f"LibreOffice conversion failed (code {process.returncode}). Check MCP server logs for details.")
            else:
                 print("INFO: LibreOffice conversion process completed.", file=sys.stderr)
                 if process.stderr: # Often has warnings even on success
                     print(f"WARNING: LibreOffice stderr: {process.stderr}", file=sys.stderr)

        except FileNotFoundError:
            print(f"ERROR: '{soffice}' command not found. Ensure LibreOffice is installed and in PATH.", file=sys.stderr)
            raise RuntimeError("LibreOffice command failed: executable not found.")
        except subprocess.TimeoutExpired:
             print(f"ERROR: LibreOffice conversion timed out after {timeout_seconds} seconds.", file=sys.stderr)
             raise RuntimeError("LibreOffice conversion timed out.")
        except Exception as e:
             print(f"ERROR: An unexpected error occurred during LibreOffice execution: {e}", file=sys.stderr)
             raise RuntimeError(f"LibreOffice execution error: {e}")


        # Find the expected output file.
        # LibreOffice often names the output based on the input file name,
        # but may vary. A common pattern is just the input filename (without ext)
        # followed by the slide number (starting from 1 for the first slide!).
        # Or sometimes just the input filename if only one slide. Let's check robustly.
        base_filename = pptx_path.stem
        expected_png_filename = f"{base_filename}{idx + 1}.png" # soffice usually uses 1-based index for output
        expected_png_filename_single = f"{base_filename}.png" # Case for single-slide conversion output

        expected_png_path = temp_dir_path / expected_png_filename
        expected_png_path_single = temp_dir_path / expected_png_filename_single

        # List files to see what was actually created
        created_files = list(temp_dir_path.glob('*.png'))
        print(f"INFO: Files created in temp dir: {[f.name for f in created_files]}", file=sys.stderr)

        actual_png_path = None
        if expected_png_path.exists():
            actual_png_path = expected_png_path
        elif expected_png_path_single.exists() and len(created_files) == 1 and idx == 0:
            # If only one PNG was created and we asked for slide 0, assume it's the one
            actual_png_path = expected_png_path_single
        elif len(created_files) > idx:
             # Fallback: If soffice just numbered them sequentially without base name (less common)
             # Sort files to try and get consistent ordering (might be fragile)
             created_files.sort()
             potential_path = created_files[idx]
             print(f"WARNING: Could not find expected PNG file name, falling back to {potential_path.name} based on index.", file=sys.stderr)
             actual_png_path = potential_path
        elif created_files:
             # If some PNGs exist but not the one we expected
             print(f"WARNING: Expected PNG file {expected_png_filename} or {expected_png_filename_single} not found, but other PNGs exist.", file=sys.stderr)
             # Maybe take the first one if index is 0? Risky.
             if idx == 0:
                  actual_png_path = created_files[0]
                  print(f"WARNING: Using first found PNG: {actual_png_path.name}", file=sys.stderr)
             else:
                  raise FileNotFoundError(f"Could not determine the correct output PNG for slide {idx} in {temp_dir}. Found: {[f.name for f in created_files]}")
        else:
            raise FileNotFoundError(f"LibreOffice ran but no PNG output files were found in {temp_dir}.")


        print(f"INFO: Reading image data from: {actual_png_path}", file=sys.stderr)
        try:
            image_bytes = actual_png_path.read_bytes()
            # Return as FastMCP Image object
            return Image(data=image_bytes, format="png")
        except Exception as e:
             print(f"ERROR: Error reading PNG file {actual_png_path}: {e}", file=sys.stderr)
             raise RuntimeError(f"Failed to read generated PNG file: {e}")


@mcp.resource("pptx://{filename}/file", mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
async def get_pptx_file(filename: str):
    """
    Provides the .pptx file for download as a binary resource.
    """
    path = _get_presentation_path(filename).resolve()
    if not path.exists():
        raise FileNotFoundError(f"Presentation file '{filename}' not found.")
    return FileResource(
        uri=f"pptx://{filename}/file",
        path=path,
        is_binary=True,
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        name=f"PPTX file: {filename}",
        description=f"Download the PowerPoint file '{filename}'."
    )


@mcp.tool()
def get_presentation_file_b64(filename: str) -> str:
    """
    Returns the content of the presentation file encoded as a Base64 string.
    Useful if the client cannot handle MCP resources or direct file paths.

    Returns:
        A Base64 encoded string representing the binary content of the .pptx file.
    """
    path = _get_presentation_path(filename)
    if not path.exists():
        raise FileNotFoundError(f"Presentation file '{filename}' not found on the server.")

    try:
        file_bytes = path.read_bytes()
        base64_encoded_bytes = base64.b64encode(file_bytes)
        base64_string = base64_encoded_bytes.decode('utf-8') # Decode bytes to string for JSON compatibility
        return base64_string
    except Exception as e:
        raise IOError(f"Error reading or encoding file '{filename}': {e}")


# --- MCP Prompts (same as before) ---
@mcp.prompt()
def flowchart_guidance() -> str:
    """Provides tips on how to create flowcharts using the available tools."""
    return """
    To create a flowchart:
    1. Use `create_or_clear_presentation` to start.
    2. Use `add_slide` with a blank layout (e.g., layout_index=6).
    3. Use `add_shape` repeatedly for flowchart elements (e.g., FLOWCHART_PROCESS, FLOWCHART_DECISION, FLOWCHART_TERMINATOR). Use `text` parameter for labels. Position using `left_inches`, `top_inches`.
    4. Use `add_shape` with connector shapes (e.g., `RIGHT_ARROW`, or find specific connectors) OR use `add_connector` tool if available (not implemented here) to connect elements. Precise positioning takes care.
    5. Check arrangement using the description resource: `pptx://{filename}/slide/{slide_index}/description` (Reliable). # Keep description as resource
    6. Optionally, render the slide image using the tool: `get_slide_image` (Requires LibreOffice setup, may be slow). # Updated to tool name
    """

@mcp.prompt()
def available_shapes() -> str:
    """Lists some common MSO_SHAPE names usable with the add_shape tool."""
    common_shapes = [
        "RECTANGLE", "OVAL", "ROUNDED_RECTANGLE", "DIAMOND", "ISOSCELES_TRIANGLE",
        "RIGHT_ARROW", "LEFT_ARROW", "UP_ARROW", "DOWN_ARROW", "PENTAGON", "HEXAGON",
        "CHEVRON", "STAR_5_POINT", "LINE_CALLOUT_1", "ACTION_BUTTON_BACK_OR_PREVIOUS",
        # Flowchart shapes
        "FLOWCHART_PROCESS", "FLOWCHART_ALTERNATE_PROCESS", "FLOWCHART_DECISION",
        "FLOWCHART_DATA", "FLOWCHART_PREDEFINED_PROCESS", "FLOWCHART_INTERNAL_STORAGE",
        "FLOWCHART_DOCUMENT", "FLOWCHART_MULTIDOCUMENT", "FLOWCHART_TERMINATOR",
        "FLOWCHART_PREPARATION", "FLOWCHART_MANUAL_INPUT", "FLOWCHART_MANUAL_OPERATION",
        "FLOWCHART_CONNECTOR", "FLOWCHART_OFFPAGE_CONNECTOR", "FLOWCHART_CARD",
        "FLOWCHART_PUNCHED_TAPE", "FLOWCHART_SUMMING_JUNCTION", "FLOWCHART_OR",
        "FLOWCHART_COLLATE", "FLOWCHART_SORT", "FLOWCHART_EXTRACT", "FLOWCHART_MERGE",
        "FLOWCHART_STORED_DATA", "FLOWCHART_DELAY", "FLOWCHART_SEQUENTIAL_ACCESS_STORAGE",
        "FLOWCHART_MAGNETIC_DISK", "FLOWCHART_DIRECT_ACCESS_STORAGE", "FLOWCHART_DISPLAY"
    ]
    return f"Common shape names for `add_shape`: {', '.join(common_shapes)}. Many others exist."


# --- Running the Server ---
if __name__ == "__main__":
    print(f"💾 Presentations will be saved in: {SAVE_DIR.resolve()}")
    print("-" * 30)
    try:
        soffice_path = _find_soffice()
        print(f"✅ Found LibreOffice executable: {soffice_path}")
    except Exception as e:
        print("❌ Image rendering (`/image.png` resource) requires LibreOffice.")
        print(f"❌ {e}")
        print("   Image rendering resource will likely fail.")
    print("-" * 30)
    print(f"🚀 Starting FastMCP server for PowerPoint generation on {HOST}:{PORT}...")
    mcp.run(transport="sse", host=HOST, port=PORT)
